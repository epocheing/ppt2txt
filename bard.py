import collections.abc
import os

import requests
from bardapi import Bard
from pptx import Presentation
from PyPDF2 import PdfReader


def read_files(input_path: str, output_path: str):
    name_list = os.listdir(input_path)
    return name_list


def ppt2text(input_path: str, name: str):
    prs = Presentation(f"{input_path}/{name}")

    prs_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                prs_text = prs_text + paragraph.text

    return prs_text


def pdf2text(input_path: str, name: str):
    reader = PdfReader(f"{input_path}/{name}")
    pages = reader.pages

    pdf_text = ""

    for page in pages:
        sub = page.extract_text()
        pdf_text = pdf_text + sub

    return pdf_text


def text2bard(TOKEN, text: str, answer: str):
    os.environ["_BARD_API_KEY"] = TOKEN

    session = requests.Session()

    session.headers = {
        "Host": "bard.google.com",
        "X-Same-Domain": "1",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.114 Safari/537.36",
        "Content-Type": "application/x-www-form-urlencoded;charset=UTF-8",
        "Origin": "https://bard.google.com",
        "Referer": "https://bard.google.com/",
    }
    session.cookies.set("__Secure-1PSID", os.getenv("_BARD_API_KEY"))

    bard = Bard(session=session, timeout=20)

    text_split = [text[i : i + 3500] for i in range(0, len(text), 3500)]

    for s in text_split:
        bard.get_answer(s)
        if s == text_split[-1]:
            answer = s + "\n" + answer
            respone = bard.get_answer(answer)["content"]

    return respone


def write_text(name: str, write_text: str, output_path: str):
    name = name.split(".")[0]
    with open(f"{output_path}/{name}.txt", "w", encoding="utf-8") as f:
        f.writelines(write_text)


def run():
    TOKEN = input("TOKEN: ")
    input_path = "./input"
    output_path = "./output"
    answer = "지금까지의 논문을 한국어로 자세하게 요약해줘"

    name_list = read_files(input_path, output_path)

    for name in name_list:
        if name.split(".")[-1] == "pdf":
            text = pdf2text(input_path, name)
            respon = text2bard(TOKEN, text, answer)
            write_text(name, respon, output_path)
        elif name.split(".")[-1] == "pptx":
            text = ppt2text(input_path, name)
            respon = text2bard(TOKEN, text, answer)
            write_text(name, respon, output_path)
        else:
            respon = "error"
            write_text(name, respon, output_path)

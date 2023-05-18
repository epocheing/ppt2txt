import collections.abc
import os

from pptx import Presentation


def ppt2txt(input_dir: str, output_dir: str):
    """pptx to txt file

    Args:
        input_dir (str): pptx dir
        output_dir (str): txt file dir
    """
    # read file names in input dir
    prs_list = os.listdir(input_dir)

    for prs_name in prs_list:
        # read pptx files
        prs = Presentation(f"./{input_dir}/{prs_name}")

        prs_txt = []

        # read pptx text
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    prs_txt.append(paragraph.text)

        # write to txt file
        with open(
            f"./{output_dir}/{prs_name.split('.')[0]}.txt", "w", encoding="utf-8"
        ) as f:
            f.writelines(prs_txt)


if __name__ == "__main__":
    ppt2txt("input", "output")
